import nltk
try:
    nltk.data.find("tokenizers/punkt")
    nltk.data.find("tokenizers/punkt_tab")
except LookupError:
    nltk.download("punkt")
    nltk.download("punkt_tab")
import streamlit as st
import fitz  # PyMuPDF for PDFs
from pptx import Presentation
import random
import re
import nltk
from nltk.corpus import stopwords
from nltk.tokenize import word_tokenize, sent_tokenize
import io

# Make sure nltk data is downloaded (only first time)
nltk.download('punkt')
nltk.download('stopwords')

# ==============================
# 📘 EXTRACT TEXT FROM FILE
# ==============================
def extract_text_from_file(uploaded_file):
    text = ""

    if uploaded_file.name.endswith(".pdf"):
        pdf_bytes = uploaded_file.read()
        with fitz.open(stream=pdf_bytes, filetype="pdf") as doc:
            for page in doc:
                text += page.get_text()

    elif uploaded_file.name.endswith((".pptx", ".ppt")):
        ppt = Presentation(uploaded_file)
        for slide in ppt.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"

    elif uploaded_file.name.endswith((".txt", ".md")):
        text = uploaded_file.read().decode("utf-8")

    else:
        st.warning("Unsupported file type. Please upload PDF, PPTX, or TXT.")
    
    return text

# ==============================
# 🧩 OFFLINE MCQ GENERATION LOGIC
# ==============================
def generate_mcqs_offline(text, num_questions=5):
    sentences = sent_tokenize(text)
    stop_words = set(stopwords.words('english'))

    # Pick meaningful sentences
    valid_sentences = [s for s in sentences if len(s.split()) > 5]
    if len(valid_sentences) < num_questions:
        num_questions = len(valid_sentences)

    selected_sentences = random.sample(valid_sentences, num_questions)
    mcqs = []

    for sent in selected_sentences:
        words = word_tokenize(sent)
        # pick a keyword (non-stopword)
        keywords = [w for w in words if w.isalpha() and w.lower() not in stop_words]
        if not keywords:
            continue
        answer = random.choice(keywords)
        question = sent.replace(answer, "______")

        # Generate options (random distractors)
        all_words = list(set(word_tokenize(text)))
        distractors = random.sample([w for w in all_words if w.isalpha() and w.lower() != answer.lower()][:100], 3)
        options = distractors + [answer]
        random.shuffle(options)

        mcq = {
            "question": question,
            "options": options,
            "answer": answer
        }
        mcqs.append(mcq)

    return mcqs

# ==============================
# 🎨 STREAMLIT UI
# ==============================
st.title("📝studybuddy")
st.markdown("Generate basic multiple-choice questions without using any API key!")

uploaded_file = st.file_uploader("📂 Upload a file", type=["pdf", "pptx", "ppt", "txt", "md"])
num_qs = st.number_input("🔢 Number of questions", min_value=1, max_value=20, value=5)

if uploaded_file is not None:
    with st.spinner("Extracting text..."):
        content = extract_text_from_file(uploaded_file)

    if st.button("✨ Generate MCQs"):
        with st.spinner("Generating MCQs... please wait ⏳"):
            mcqs = generate_mcqs_offline(content, num_qs)

        if not mcqs:
            st.warning("Not enough text to generate MCQs.")
        else:
            st.success("✅ MCQs Generated Successfully!")

            output_text = ""
            for i, q in enumerate(mcqs, 1):
                st.markdown(f"**Q{i}. {q['question']}**")
                for opt in q["options"]:
                    st.write(f"- {opt}")
                st.write(f"**Answer:** {q['answer']}")
                st.write("---")

                output_text += f"Q{i}. {q['question']}\n"
                for opt in q["options"]:
                    output_text += f"- {opt}\n"
                output_text += f"Answer: {q['answer']}\n\n"

            st.download_button(
                label="📥 Download MCQs as TXT",
                data=output_text.encode("utf-8"),
                file_name="offline_mcqs.txt",
                mime="text/plain"
            )
